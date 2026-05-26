from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import pickle
import os
import pypdf
import docx
from pptx import Presentation
import ollama
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Lazy load model to save memory on startup
_model = None
documents = []
index = None
memory = []
MAX_MEMORY = 20  # Keep only last 20 messages


def get_model():
    """Lazy load the embedding model on first use"""
    global _model
    if _model is None:
        logger.info("Loading embedding model...")
        _model = SentenceTransformer("all-MiniLM-L6-v2")
        logger.info("✅ Model loaded")
    return _model


def extract_text(file):
    """Extract text from various file formats"""
    try:
        if file.endswith(".pdf"):
            reader = pypdf.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text

        if file.endswith(".txt"):
            with open(file, "r", encoding="utf-8") as f:
                return f.read()

        if file.endswith(".docx"):
            doc = docx.Document(file)
            return "\n".join([p.text for p in doc.paragraphs])

        if file.endswith(".pptx"):
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        return ""
    except Exception as e:
        logger.error(f"Error extracting text: {e}")
        return ""


def chunk_text(text, size=400, overlap=50):
    """Split text into overlapping chunks for better context"""
    chunks = []
    
    # Clean text
    text = text.strip()
    if not text:
        return chunks
    
    for i in range(0, len(text), size - overlap):
        chunk = text[i:i+size].strip()
        if len(chunk) > 50:  # Only keep meaningful chunks
            chunks.append(chunk)
    
    return chunks


def create_vector_store(chunks):
    """Create or update FAISS vector store"""
    global index, documents

    try:
        if not chunks:
            logger.warning("No chunks to index")
            return

        documents.extend(chunks)
        logger.info(f"Encoding {len(chunks)} chunks...")
        
        embeddings = get_model().encode(chunks, convert_to_numpy=True)

        if index is None:
            dimension = embeddings.shape[1]
            index = faiss.IndexFlatL2(dimension)
            logger.info(f"Created new index with dimension {dimension}")

        index.add(np.array(embeddings, dtype=np.float32))
        logger.info(f"Index now has {index.ntotal} vectors")

        # Save to disk
        os.makedirs("faiss_index", exist_ok=True)
        faiss.write_index(index, "faiss_index/index.faiss")

        with open("faiss_index/chunks.pkl", "wb") as f:
            pickle.dump(documents, f)
        
        logger.info("✅ Vector store saved")

    except Exception as e:
        logger.error(f"Error creating vector store: {e}")
        raise


def load_vector_store():
    """Load FAISS index from disk"""
    global index, documents

    try:
        if os.path.exists("faiss_index/index.faiss"):
            logger.info("Loading vector store from disk...")
            index = faiss.read_index("faiss_index/index.faiss")

            with open("faiss_index/chunks.pkl", "rb") as f:
                documents = pickle.load(f)
            
            logger.info(f"✅ Loaded index with {index.ntotal} vectors")
        else:
            logger.info("No existing vector store found")
    
    except Exception as e:
        logger.error(f"Error loading vector store: {e}")


def process_file(path):
    """Process uploaded file and add to vector store"""
    try:
        logger.info(f"Processing file: {path}")
        text = extract_text(path)

        if not text or len(text.strip()) < 50:
            raise ValueError("File contains insufficient text")

        logger.info(f"Extracted {len(text)} characters")
        chunks = chunk_text(text)
        logger.info(f"Created {len(chunks)} chunks")

        create_vector_store(chunks)
        logger.info("✅ File processed successfully")

    except Exception as e:
        logger.error(f"Error processing file: {e}")
        raise


def add_to_memory(question, answer):
    """Add conversation to memory with limit"""
    global memory
    memory.append({"role": "user", "content": question})
    memory.append({"role": "assistant", "content": answer})
    
    # Keep only last MAX_MEMORY messages
    if len(memory) > MAX_MEMORY:
        memory = memory[-MAX_MEMORY:]


def format_memory():
    """Format memory for context"""
    if not memory:
        return ""
    
    history = "\n".join([
        f"{msg['role'].title()}: {msg['content'][:100]}..."
        for msg in memory[-4:]  # Last 4 messages
    ])
    return history


def ask_question(question):
    """Answer question based on uploaded documents"""
    try:
        global memory

        if index is None or len(documents) == 0:
            return "❌ No documents uploaded yet. Please upload a file first."

        logger.info(f"Searching for context...")
        query_embedding = get_model().encode([question], convert_to_numpy=True)

        # Search for top 3 relevant chunks
        distances, indices = index.search(np.array([query_embedding], dtype=np.float32), k=3)

        context = "\n---\n".join([
            documents[i] for i in indices[0] if i < len(documents)
        ])

        memory_context = format_memory()

        prompt = f"""You are a helpful AI assistant analyzing documents.
        
Based on the following document excerpts, answer the user's question clearly and concisely.

DOCUMENT CONTEXT:
{context}

PREVIOUS CONVERSATION:
{memory_context if memory_context else "None"}

USER QUESTION:
{question}

ANSWER:"""

        logger.info("Calling Ollama...")
        response = ollama.chat(
            model="phi3",
            messages=[{"role": "user", "content": prompt}],
            stream=False
        )

        answer = response["message"]["content"].strip()
        add_to_memory(question, answer)
        
        logger.info("✅ Question answered")
        return answer

    except Exception as e:
        logger.error(f"Error in ask_question: {e}")
        raise


def general_ai(question):
    """Answer general AI questions without documents"""
    try:
        global memory

        memory_context = format_memory()

        prompt = f"""You are a helpful, friendly AI assistant.
        
PREVIOUS CONVERSATION:
{memory_context if memory_context else "None"}

USER QUESTION:
{question}

ANSWER:"""

        logger.info("Calling Ollama for general AI...")
        response = ollama.chat(
            model="phi3",
            messages=[{"role": "user", "content": prompt}],
            stream=False
        )

        answer = response["message"]["content"].strip()
        add_to_memory(question, answer)
        
        logger.info("✅ General AI response generated")
        return answer

    except Exception as e:
        logger.error(f"Error in general_ai: {e}")
        raise